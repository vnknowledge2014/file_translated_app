"""Tests for deterministic document reconstruction."""

import json
import os
import tempfile

import pytest

from app.agent.reconstructor import (
    _build_translation_map,
    _replace_in_text,
    reconstruct_docx,
    reconstruct_xlsx,
    reconstruct_pptx,
    reconstruct_plaintext,
    reconstruct_document,
)


class TestBuildTranslationMap:
    def test_basic(self):
        segments = [
            {"text": "日本語", "translated_text": "Tiếng Nhật"},
            {"text": "テスト", "translated_text": "Kiểm tra"},
        ]
        tmap = _build_translation_map(segments)
        assert tmap == {"日本語": "Tiếng Nhật", "テスト": "Kiểm tra"}

    def test_skips_empty(self):
        segments = [
            {"text": "", "translated_text": "Hello"},
            {"text": "JP", "translated_text": ""},
            {"text": "Same", "translated_text": "Same"},  # same → skipped
        ]
        tmap = _build_translation_map(segments)
        assert tmap == {}

    def test_strips_whitespace(self):
        segments = [
            {"text": "  テスト  ", "translated_text": "  Kiểm tra  "},
        ]
        tmap = _build_translation_map(segments)
        assert tmap == {"テスト": "Kiểm tra"}


class TestReplaceInText:
    def test_exact_match(self):
        tmap = {"日本語": "Tiếng Nhật"}
        assert _replace_in_text("日本語", tmap) == "Tiếng Nhật"

    def test_no_match(self):
        tmap = {"日本語": "Tiếng Nhật"}
        assert _replace_in_text("English", tmap) is None

    def test_empty_text(self):
        tmap = {"日本語": "Tiếng Nhật"}
        assert _replace_in_text("", tmap) is None
        assert _replace_in_text("   ", tmap) is None

    def test_partial_match(self):
        tmap = {"日本語": "Tiếng Nhật"}
        result = _replace_in_text("Test 日本語 here", tmap)
        assert result == "Test Tiếng Nhật here"


class TestReconstructDocx:
    def test_basic_replacement(self):
        """Test DOCX reconstruction with a real .docx file."""
        from docx import Document

        with tempfile.TemporaryDirectory() as tmpdir:
            # Create a simple DOCX
            doc = Document()
            p = doc.add_paragraph()
            run = p.add_run("テスト文書")
            run.bold = True
            doc.add_paragraph("普通のテキスト")
            src = os.path.join(tmpdir, "test.docx")
            doc.save(src)

            # Translate
            segments = [
                {"text": "テスト文書", "translated_text": "Tài liệu thử nghiệm"},
                {"text": "普通のテキスト", "translated_text": "Văn bản thông thường"},
            ]
            out = os.path.join(tmpdir, "test_vi.docx")
            result = reconstruct_docx(src, segments, out)

            # Verify
            assert os.path.exists(result)
            doc2 = Document(result)
            assert doc2.paragraphs[0].text == "Tài liệu thử nghiệm"
            assert doc2.paragraphs[1].text == "Văn bản thông thường"
            # Verify bold is preserved
            assert doc2.paragraphs[0].runs[0].bold is True

    def test_table_replacement(self):
        """Test DOCX table cell reconstruction."""
        from docx import Document

        with tempfile.TemporaryDirectory() as tmpdir:
            doc = Document()
            table = doc.add_table(rows=2, cols=2)
            table.cell(0, 0).text = "名前"
            table.cell(0, 1).text = "値"
            table.cell(1, 0).text = "テスト"
            table.cell(1, 1).text = "100"
            src = os.path.join(tmpdir, "table.docx")
            doc.save(src)

            segments = [
                {"text": "名前", "translated_text": "Tên"},
                {"text": "値", "translated_text": "Giá trị"},
                {"text": "テスト", "translated_text": "Kiểm tra"},
            ]
            out = os.path.join(tmpdir, "table_vi.docx")
            reconstruct_docx(src, segments, out)

            doc2 = Document(out)
            assert doc2.tables[0].cell(0, 0).text == "Tên"
            assert doc2.tables[0].cell(0, 1).text == "Giá trị"
            assert doc2.tables[0].cell(1, 0).text == "Kiểm tra"
            assert doc2.tables[0].cell(1, 1).text == "100"  # untranslated

    def test_empty_segments(self):
        """No segments → just copy the file."""
        from docx import Document

        with tempfile.TemporaryDirectory() as tmpdir:
            doc = Document()
            doc.add_paragraph("テスト")
            src = os.path.join(tmpdir, "orig.docx")
            doc.save(src)

            out = os.path.join(tmpdir, "out.docx")
            reconstruct_docx(src, [], out)
            assert os.path.exists(out)
            # Text unchanged
            doc2 = Document(out)
            assert doc2.paragraphs[0].text == "テスト"


class TestReconstructXlsx:
    def test_basic_replacement(self):
        """Test XLSX cell value replacement."""
        from openpyxl import Workbook, load_workbook

        with tempfile.TemporaryDirectory() as tmpdir:
            wb = Workbook()
            ws = wb.active
            ws["A1"] = "項目名"
            ws["B1"] = "説明"
            ws["A2"] = "テスト"
            ws["B2"] = 42  # number — should be untouched
            # Add some styling
            ws["A1"].font = ws["A1"].font.copy(bold=True)
            src = os.path.join(tmpdir, "test.xlsx")
            wb.save(src)

            segments = [
                {"text": "項目名", "translated_text": "Tên mục"},
                {"text": "説明", "translated_text": "Mô tả"},
                {"text": "テスト", "translated_text": "Kiểm tra"},
            ]
            out = os.path.join(tmpdir, "test_vi.xlsx")
            reconstruct_xlsx(src, segments, out)

            wb2 = load_workbook(out)
            ws2 = wb2.active
            assert ws2["A1"].value == "Tên mục"
            assert ws2["B1"].value == "Mô tả"
            assert ws2["A2"].value == "Kiểm tra"
            assert ws2["B2"].value == 42  # number preserved
            assert ws2["A1"].font.bold is True  # style preserved

    def test_sheet_name_rename(self):
        """Sheet names should be translated during reconstruction."""
        from openpyxl import Workbook, load_workbook

        with tempfile.TemporaryDirectory() as tmpdir:
            wb = Workbook()
            ws = wb.active
            ws.title = "API作成_スケジュール"
            ws["A1"] = "テスト"
            src = os.path.join(tmpdir, "test.xlsx")
            wb.save(src)

            segments = [
                {"text": "テスト", "translated_text": "Kiểm tra"},
                {"text": "API作成_スケジュール", "translated_text": "Lịch trình tạo API", "type": "sheet_name"},
            ]
            out = os.path.join(tmpdir, "test_vi.xlsx")
            reconstruct_xlsx(src, segments, out)

            wb2 = load_workbook(out)
            assert wb2.sheetnames[0] == "Lịch trình tạo API"
            assert wb2.active["A1"].value == "Kiểm tra"


class TestReconstructPptx:
    def test_basic_replacement(self):
        """Test PPTX text frame replacement."""
        from pptx import Presentation
        from pptx.util import Inches

        with tempfile.TemporaryDirectory() as tmpdir:
            prs = Presentation()
            slide = prs.slides.add_slide(prs.slide_layouts[1])  # title + content
            title = slide.shapes.title
            title.text = "テストスライド"
            body = slide.placeholders[1]
            body.text = "内容テキスト"
            src = os.path.join(tmpdir, "test.pptx")
            prs.save(src)

            segments = [
                {"text": "テストスライド", "translated_text": "Slide thử nghiệm"},
                {"text": "内容テキスト", "translated_text": "Nội dung văn bản"},
            ]
            out = os.path.join(tmpdir, "test_vi.pptx")
            reconstruct_pptx(src, segments, out)

            prs2 = Presentation(out)
            slide2 = prs2.slides[0]
            assert slide2.shapes.title.text == "Slide thử nghiệm"


class TestReconstructPlaintext:
    def test_basic_lines(self):
        with tempfile.TemporaryDirectory() as tmpdir:
            src = os.path.join(tmpdir, "test.txt")
            with open(src, "w") as f:
                f.write("日本語テキスト\n")
                f.write("English line\n")
                f.write("テスト\n")

            segments = [
                {"text": "日本語テキスト", "translated_text": "Văn bản tiếng Nhật", "location": "line[0]"},
                {"text": "テスト", "translated_text": "Kiểm tra", "location": "line[2]"},
            ]
            out = os.path.join(tmpdir, "test_vi.txt")
            reconstruct_plaintext(src, segments, out)

            with open(out) as f:
                lines = f.readlines()
            assert lines[0].strip() == "Văn bản tiếng Nhật"
            assert lines[1].strip() == "English line"
            assert lines[2].strip() == "Kiểm tra"

    def test_markdown_prefix(self):
        with tempfile.TemporaryDirectory() as tmpdir:
            src = os.path.join(tmpdir, "test.md")
            with open(src, "w") as f:
                f.write("## タイトル\n")
                f.write("- 項目一\n")

            segments = [
                {"text": "## タイトル", "translated_text": "Tiêu đề", "location": "line[0]"},
                {"text": "- 項目一", "translated_text": "Mục một", "location": "line[1]"},
            ]
            out = os.path.join(tmpdir, "test_vi.md")
            reconstruct_plaintext(src, segments, out)

            with open(out) as f:
                lines = f.readlines()
            assert lines[0].strip() == "## Tiêu đề"
            assert lines[1].strip() == "- Mục một"

    def test_diagram_token_reconstruction(self):
        """Diagram tokens should replace JP text in-place, preserving box structure."""
        with tempfile.TemporaryDirectory() as tmpdir:
            src = os.path.join(tmpdir, "diag.md")
            with open(src, "w") as f:
                f.write("```\n")
                f.write("┌──────────┐ ┌──────────┐\n")
                f.write("│ サービス │ │ ナレッジ │\n")
                f.write("└──────────┘ └──────────┘\n")
                f.write("```\n")

            segments = [
                {"text": "サービス", "translated_text": "Dịch vụ", "location": "line[2]", "type": "diagram_token"},
                {"text": "ナレッジ", "translated_text": "Kiến thức", "location": "line[2]", "type": "diagram_token"},
            ]
            out = os.path.join(tmpdir, "diag_vi.md")
            reconstruct_plaintext(src, segments, out)

            with open(out) as f:
                lines = f.readlines()
            # Both translations fit within available space (no grid expansion needed):
            # "Dịch vụ" (vw=7) in available=10 → padded with 3 spaces
            # "Kiến thức" (vw=9) in available=10 → padded with 1 space
            assert "Dịch vụ" in lines[2]
            assert "Kiến thức" in lines[2]
            # Border lines preserved
            assert "┌──────────┐" in lines[1]


class TestReconstructDispatcher:
    def test_unsupported_type(self):
        with pytest.raises(ValueError, match="No deterministic"):
            reconstruct_document("pdf", "/fake", [], "/fake/out")

    def test_dispatches_docx(self):
        from docx import Document

        with tempfile.TemporaryDirectory() as tmpdir:
            doc = Document()
            doc.add_paragraph("テスト")
            src = os.path.join(tmpdir, "test.docx")
            doc.save(src)

            segments = [
                {"text": "テスト", "translated_text": "Kiểm tra"},
            ]
            out = os.path.join(tmpdir, "test_vi.docx")
            result = reconstruct_document("docx", src, segments, out)
            assert os.path.exists(result)
