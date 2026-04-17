"""Tests for deterministic document extraction."""

import os
import tempfile

import pytest

from app.agent.extractor import (
    _is_translatable,
    _dedup_segments,
    extract_docx,
    extract_xlsx,
    extract_pptx,
    extract_plaintext,
    extract_document,
)


class TestIsTranslatable:
    def test_japanese_text(self):
        assert _is_translatable("テスト") is True
        assert _is_translatable("日本語テキスト") is True

    def test_english_text(self):
        assert _is_translatable("Hello world") is False

    def test_empty_text(self):
        assert _is_translatable("") is False
        assert _is_translatable(None) is False
        assert _is_translatable("   ") is False

    def test_url_with_jp(self):
        assert _is_translatable("https://example.com/テスト") is False

    def test_email(self):
        assert _is_translatable("test@example.com") is False

    def test_formula(self):
        assert _is_translatable("=SUM(A1:A10)") is False

    def test_pure_numbers(self):
        assert _is_translatable("12345") is False
        assert _is_translatable("99.9%") is False

    def test_mixed_jp_en(self):
        assert _is_translatable("Hello テスト world") is True


class TestDedupSegments:
    def test_no_dups(self):
        segs = [{"text": "A"}, {"text": "B"}]
        assert len(_dedup_segments(segs)) == 2

    def test_with_dups(self):
        segs = [{"text": "テスト"}, {"text": "別の"}, {"text": "テスト"}]
        result = _dedup_segments(segs)
        assert len(result) == 2
        assert result[0]["text"] == "テスト"
        assert result[1]["text"] == "別の"


class TestExtractDocx:
    def test_basic_paragraphs(self):
        from docx import Document

        with tempfile.TemporaryDirectory() as tmpdir:
            doc = Document()
            doc.add_paragraph("テスト文書")
            doc.add_paragraph("普通のテキスト")
            doc.add_paragraph("English only")  # Should be skipped
            path = os.path.join(tmpdir, "test.docx")
            doc.save(path)

            segs = extract_docx(path)
            assert len(segs) == 2
            assert segs[0]["text"] == "テスト文書"
            assert segs[1]["text"] == "普通のテキスト"

    def test_table_extraction(self):
        from docx import Document

        with tempfile.TemporaryDirectory() as tmpdir:
            doc = Document()
            table = doc.add_table(rows=2, cols=2)
            table.cell(0, 0).text = "名前"
            table.cell(0, 1).text = "値"
            table.cell(1, 0).text = "English"
            table.cell(1, 1).text = "テスト"
            path = os.path.join(tmpdir, "table.docx")
            doc.save(path)

            segs = extract_docx(path)
            texts = {s["text"] for s in segs}
            assert "名前" in texts
            assert "値" in texts
            assert "テスト" in texts
            assert "English" not in texts  # No Japanese

    def test_empty_doc(self):
        from docx import Document

        with tempfile.TemporaryDirectory() as tmpdir:
            doc = Document()
            doc.add_paragraph("")
            path = os.path.join(tmpdir, "empty.docx")
            doc.save(path)

            segs = extract_docx(path)
            assert len(segs) == 0

    def test_dedup(self):
        """Same text in multiple cells → only extract once."""
        from docx import Document

        with tempfile.TemporaryDirectory() as tmpdir:
            doc = Document()
            doc.add_paragraph("テスト")
            table = doc.add_table(rows=1, cols=1)
            table.cell(0, 0).text = "テスト"  # Same text as paragraph
            path = os.path.join(tmpdir, "dup.docx")
            doc.save(path)

            segs = extract_docx(path)
            assert len(segs) == 1


class TestExtractXlsx:
    def test_basic_cells(self):
        from openpyxl import Workbook

        with tempfile.TemporaryDirectory() as tmpdir:
            wb = Workbook()
            ws = wb.active
            ws["A1"] = "項目名"
            ws["B1"] = "説明"
            ws["A2"] = 42          # Number — skip
            ws["B2"] = "テスト"
            ws["C2"] = "English"   # No JP — skip
            path = os.path.join(tmpdir, "test.xlsx")
            wb.save(path)

            segs = extract_xlsx(path)
            texts = {s["text"] for s in segs}
            assert "項目名" in texts
            assert "説明" in texts
            assert "テスト" in texts
            assert "English" not in texts

    def test_multiple_sheets(self):
        from openpyxl import Workbook

        with tempfile.TemporaryDirectory() as tmpdir:
            wb = Workbook()
            ws1 = wb.active
            ws1.title = "Sheet1"
            ws1["A1"] = "シート1"

            ws2 = wb.create_sheet("Sheet2")
            ws2["A1"] = "シート2"

            path = os.path.join(tmpdir, "multi.xlsx")
            wb.save(path)

            segs = extract_xlsx(path)
            cell_segs = [s for s in segs if s["type"] != "sheet_name"]
            assert len(cell_segs) == 2

    def test_sheet_name_extraction(self):
        """Sheet names with Japanese text should be extracted."""
        from openpyxl import Workbook

        with tempfile.TemporaryDirectory() as tmpdir:
            wb = Workbook()
            ws1 = wb.active
            ws1.title = "API作成_スケジュール"
            ws1["A1"] = "テスト"

            ws2 = wb.create_sheet("EnglishSheet")
            ws2["A1"] = "テスト2"

            path = os.path.join(tmpdir, "sheets.xlsx")
            wb.save(path)

            segs = extract_xlsx(path)
            sheet_segs = [s for s in segs if s["type"] == "sheet_name"]
            assert len(sheet_segs) == 1
            assert sheet_segs[0]["text"] == "API作成_スケジュール"
            assert sheet_segs[0]["location"] == "sheet_name:API作成_スケジュール"


class TestExtractPptx:
    def test_basic_slides(self):
        from pptx import Presentation

        with tempfile.TemporaryDirectory() as tmpdir:
            prs = Presentation()
            slide = prs.slides.add_slide(prs.slide_layouts[1])
            title = slide.shapes.title
            title.text = "テストスライド"
            body = slide.placeholders[1]
            body.text = "内容テキスト"
            path = os.path.join(tmpdir, "test.pptx")
            prs.save(path)

            segs = extract_pptx(path)
            texts = {s["text"] for s in segs}
            assert "テストスライド" in texts
            assert "内容テキスト" in texts

    def test_multiple_slides(self):
        from pptx import Presentation

        with tempfile.TemporaryDirectory() as tmpdir:
            prs = Presentation()
            for i in range(3):
                slide = prs.slides.add_slide(prs.slide_layouts[0])
                slide.shapes.title.text = f"スライド{i+1}"
            path = os.path.join(tmpdir, "multi.pptx")
            prs.save(path)

            segs = extract_pptx(path)
            assert len(segs) == 3


class TestExtractPlaintext:
    def test_basic_lines(self):
        with tempfile.TemporaryDirectory() as tmpdir:
            path = os.path.join(tmpdir, "test.txt")
            with open(path, "w") as f:
                f.write("日本語テキスト\n")
                f.write("English only\n")
                f.write("テスト二行目\n")

            segs = extract_plaintext(path)
            assert len(segs) == 2
            assert segs[0]["text"] == "日本語テキスト"
            assert segs[0]["location"] == "line[0]"
            assert segs[1]["location"] == "line[2]"

    def test_code_block_skip(self):
        with tempfile.TemporaryDirectory() as tmpdir:
            path = os.path.join(tmpdir, "test.md")
            with open(path, "w") as f:
                f.write("# タイトル\n")
                f.write("```python\n")
                f.write("# これはコードです\n")  # Inside code block — skip
                f.write("```\n")
                f.write("テスト\n")

            segs = extract_plaintext(path)
            texts = [s["text"] for s in segs]
            assert "# タイトル" in texts
            assert "テスト" in texts
            assert "# これはコードです" not in texts

    def test_preserves_duplicates(self):
        """Plaintext should NOT dedup (location-based reconstruction)."""
        with tempfile.TemporaryDirectory() as tmpdir:
            path = os.path.join(tmpdir, "test.txt")
            with open(path, "w") as f:
                f.write("テスト\n")
                f.write("テスト\n")

            segs = extract_plaintext(path)
            assert len(segs) == 2  # Not deduped

    def test_diagram_code_block_extraction(self):
        """ASCII diagram code blocks should extract JP tokens, not skip."""
        with tempfile.TemporaryDirectory() as tmpdir:
            path = os.path.join(tmpdir, "diag.md")
            with open(path, "w") as f:
                f.write("# タイトル\n")
                f.write("```\n")
                f.write("┌──────────┐\n")
                f.write("│ サービス │\n")
                f.write("└──────────┘\n")
                f.write("```\n")
                f.write("テスト\n")

            segs = extract_plaintext(path)
            texts = [s["text"] for s in segs]
            types = [s["type"] for s in segs]

            # Should have: タイトル (body), サービス (diagram_token), テスト (body)
            assert "# タイトル" in texts
            assert "サービス" in texts
            assert "テスト" in texts
            assert "diagram_token" in types

    def test_real_code_block_still_skipped(self):
        """Regular code blocks (no box chars) should still be skipped."""
        with tempfile.TemporaryDirectory() as tmpdir:
            path = os.path.join(tmpdir, "code.md")
            with open(path, "w") as f:
                f.write("```python\n")
                f.write("# これはコードです\n")
                f.write("print('テスト')\n")
                f.write("```\n")

            segs = extract_plaintext(path)
            assert len(segs) == 0  # All inside code block

    def test_diagram_multi_token_per_line(self):
        """Multiple JP tokens on one diagram line should all be extracted."""
        with tempfile.TemporaryDirectory() as tmpdir:
            path = os.path.join(tmpdir, "multi.md")
            with open(path, "w") as f:
                f.write("```\n")
                f.write("│ サービス │ ナレッジ │\n")
                f.write("```\n")

            segs = extract_plaintext(path)
            texts = [s["text"] for s in segs]
            assert "サービス" in texts
            assert "ナレッジ" in texts
            assert len(segs) == 2


class TestExtractDispatcher:
    def test_unsupported_type(self):
        with pytest.raises(ValueError, match="No deterministic"):
            extract_document("pdf", "/fake")

    def test_file_not_found(self):
        with pytest.raises(FileNotFoundError):
            extract_document("docx", "/nonexistent/file.docx")

    def test_dispatches_docx(self):
        from docx import Document

        with tempfile.TemporaryDirectory() as tmpdir:
            doc = Document()
            doc.add_paragraph("テスト")
            path = os.path.join(tmpdir, "test.docx")
            doc.save(path)

            segs = extract_document("docx", path)
            assert len(segs) == 1
